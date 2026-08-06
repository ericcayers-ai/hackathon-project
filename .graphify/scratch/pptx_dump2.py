from pptx import Presentation
p = Presentation('docs/ClearChart(nurse_notes)_Pitch.pptx')
with open('.graphify/scratch/pptx_dump2.txt', 'w', encoding='utf-8') as f:
    f.write(f"SLIDES: {len(p.slides)}\n")
    for i, s in enumerate(p.slides):
        f.write(f"--- Slide {i+1} ---\n")
        for sh in s.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text
                if t.strip():
                    f.write(t + "\n")
            if sh.shape_type == 6:  # group
                for sub in sh.shapes:
                    if sub.has_text_frame and sub.text_frame.text.strip():
                        f.write(sub.text_frame.text + "\n")
