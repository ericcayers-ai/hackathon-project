# -*- coding: utf-8 -*-
"""Add a condensed 'Technical roadmap' line to slide 8 of the ClearChart pptx,
summarizing the backend proposal now in the docx (MedGemma extraction as the
realistic next step; de-identification pipeline + simplification fine-tuning
named as future roadmap)."""
from pptx import Presentation
from pptx.util import Emu

PATH = "docs/ClearChart(nurse_notes)_Pitch.pptx"
p = Presentation(PATH)
slide8 = p.slides[7]  # 0-indexed

# Find the "What we need" text box to clone formatting from, and locate an
# empty area to place a new small textbox with the roadmap note.
target = None
for sh in slide8.shapes:
    if sh.has_text_frame and "What we need" in sh.text_frame.text:
        target = sh
        break
assert target is not None, "could not find 'What we need' shape on slide 8"

# Add a new textbox just below/right using the same left/width, offset in top.
left = target.left
top = target.top + target.height + Emu(91440)  # ~0.1 inch gap
width = target.width
height = Emu(457200)  # ~0.5 inch, will autosize with text

tb = slide8.shapes.add_textbox(left, top, width, height)
tf = tb.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
run_label = p0.add_run()
run_label.text = "Technical roadmap:  "
run_label.font.bold = True
run_label.font.size = target.text_frame.paragraphs[0].runs[0].font.size if target.text_frame.paragraphs[0].runs else None
run_body = p0.add_run()
run_body.text = ("Next build step is MedGemma 4B for structured medicine/date extraction "
                  "(targets the 6/52 misparse directly). On-device de-identification "
                  "(LLM-AIx/OpenMed) and simplification fine-tuning are scoped as post-hackathon "
                  "future work, not claimed as built.")

p.save(PATH)
print("Added technical roadmap note to slide 8.")
