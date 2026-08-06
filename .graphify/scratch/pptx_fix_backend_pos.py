# -*- coding: utf-8 -*-
"""Fix: the technical-roadmap textbox added to slide 8 overlapped the
'What we need' bullet column. Reposition it into the clear space below
all three columns instead, spanning full content width."""
from pptx import Presentation
from pptx.util import Emu

PATH = "docs/ClearChart(nurse_notes)_Pitch.pptx"
p = Presentation(PATH)
slide8 = p.slides[7]

roadmap_shape = None
for sh in slide8.shapes:
    if sh.has_text_frame and "Technical roadmap" in sh.text_frame.text:
        roadmap_shape = sh
        break
assert roadmap_shape is not None

# Bottom-most existing content bottom edge (Synthetic sample docs bullet: top 5175504 + height 566928)
bottom_of_content = 5175504 + 566928  # = 5742432
new_top = bottom_of_content + Emu(45720)  # ~0.05" gap
new_left = 822960  # align with slide title / "What changes" left margin
new_width = 10515600  # full content width, matches "What changes if this ships"

roadmap_shape.top = new_top
roadmap_shape.left = new_left
roadmap_shape.width = new_width
roadmap_shape.height = Emu(457200)

p.save(PATH)
print("Repositioned technical roadmap textbox to top=%d left=%d width=%d" % (new_top, new_left, new_width))
